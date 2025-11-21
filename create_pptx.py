from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_dark_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Dark theme colors
    bg_color = RGBColor(30, 30, 30)  # Dark background
    title_color = RGBColor(100, 200, 255)  # Light blue
    text_color = RGBColor(220, 220, 220)  # Light gray
    accent_color = RGBColor(255, 180, 0)  # Orange accent
    code_bg = RGBColor(45, 45, 45)  # Darker gray for code
    code_text = RGBColor(200, 255, 200)  # Light green for code

    def add_slide_with_title(title):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = title_color

        return slide

    def add_content_box(slide, text, top, height=5, is_code=False):
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(height))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = text

        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(14 if not is_code else 12)
            paragraph.font.color.rgb = code_text if is_code else text_color
            paragraph.font.name = 'Consolas' if is_code else 'Segoe UI'

        if is_code:
            shape = content_box
            shape.fill.solid()
            shape.fill.fore_color.rgb = code_bg
            shape.line.color.rgb = RGBColor(80, 80, 80)

        return content_box

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "📘 JAVA BASIC TRAINING"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    stf = subtitle_box.text_frame
    stf.text = "Automation Testing & Selenium Foundation"
    sp = stf.paragraphs[0]
    sp.font.size = Pt(24)
    sp.font.color.rgb = accent_color
    sp.alignment = PP_ALIGN.CENTER

    # Slide 2: Objectives
    slide2 = add_slide_with_title("Slide 1 – Java Basic Training – Giới thiệu")
    content = """🎯 Mục tiêu khóa học:
• Hiểu rõ biến, kiểu dữ liệu, mảng, điều kiện, hàm trong Java
• Làm nền tảng để học Selenium và Automation Testing

📋 Nội dung chính:
1. Biến & Kiểu dữ liệu
2. Câu lệnh điều kiện
3. Mảng 1D – 2D
4. Hàm (method)"""
    add_content_box(slide2, content, 1.3, 5.5)

    # Slide 3: Variables Introduction
    slide3 = add_slide_with_title("Slide 2 – Biến (Variables) – Giới thiệu")
    content = """📌 Định nghĩa:
Biến là vùng nhớ lưu dữ liệu trong chương trình.

💻 Cú pháp: kiểuDữLiệu tênBiến = giáTrị;

✅ Quy tắc đặt tên biến:
• ❌ Không bắt đầu bằng số
• ❌ Không dùng ký tự đặc biệt (trừ _ và $)
• ✅ Dùng camelCase cho biến & method

🔖 Các loại biến trong Java:
1. Local Variable (Biến cục bộ) – trong method/block
2. Instance Variable (Biến instance) – trong class, thuộc object
3. Static Variable (Biến tĩnh) – trong class, dùng chung"""
    add_content_box(slide3, content, 1.3, 5.5)

    # Slide 4: Local Variable
    slide4 = add_slide_with_title("Slide 3 – Local Variable (Biến cục bộ)")
    content = """📍 Khai báo: Bên trong method, constructor hoặc block { }

🔍 Phạm vi: Chỉ dùng trong method/block đó

⚡ Đặc điểm:
• Không có giá trị mặc định → Phải khởi tạo trước khi dùng
• Lưu trong Stack memory
• Tồn tại cho đến khi method kết thúc

💻 Ví dụ:
public void calculateTotal() {
    int quantity = 5;           // local variable
    double price = 100.0;       // local variable
    double total = quantity * price;
    System.out.println("Total: " + total);
}"""
    add_content_box(slide4, content, 1.3, 5.5)

    # Slide 5: Instance Variable
    slide5 = add_slide_with_title("Slide 4 – Instance Variable (Biến instance)")
    content = """📍 Khai báo: Trong class nhưng ngoài method

🔍 Phạm vi: Dùng được ở mọi method trong class (qua object)

⚡ Đặc điểm:
• Có giá trị mặc định (int=0, double=0.0, boolean=false, null)
• Lưu trong Heap memory
• Mỗi object có 1 bản sao riêng

💻 Ví dụ:
class Student {
    String name;    // instance variable
    int age;        // instance variable
    double gpa;
}

Student student1 = new Student();
student1.name = "John";"""
    add_content_box(slide5, content, 1.3, 5.5)

    # Slide 6: Static Variable
    slide6 = add_slide_with_title("Slide 5 – Static Variable (Biến tĩnh)")
    content = """📍 Khai báo: Trong class với từ khóa static

🔍 Phạm vi: Thuộc về class, được chia sẻ chung

⚡ Đặc điểm:
• Chỉ có 1 bản duy nhất trong bộ nhớ
• Có giá trị mặc định
• Lưu trong Method Area
• Tồn tại cho đến khi chương trình kết thúc

💻 Ví dụ:
public class Calculator {
    static int calculationCount = 0;

    public static void add(int a, int b) {
        calculationCount++;
        System.out.println("Total: " + calculationCount);
    }
}"""
    add_content_box(slide6, content, 1.3, 5.5)

    # Slide 7: Variable Comparison
    slide7 = add_slide_with_title("Slide 6 – So sánh 3 loại biến")
    content = """📊 Bảng so sánh:

Đặc điểm          Local          Instance       Static
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Vị trí khai báo   method/block   class,ngoài    class+static
Giá trị mặc định  ❌ Không có    ✅ Có          ✅ Có
Phạm vi           method/block   toàn object    toàn class
Bộ nhớ            Stack          Heap           Method Area
Số bản sao        1/lần gọi      1/object       1/class

🎯 Khi nào dùng:
✅ Local: Biến tạm thời trong method
✅ Instance: Thuộc tính của object (name, age...)
✅ Static: Hằng số chung, biến đếm, configuration"""
    add_content_box(slide7, content, 1.3, 5.5)

    # Slide 8: Primitive Types
    slide8 = add_slide_with_title("Slide 7 – Kiểu dữ liệu Primitive")
    content = """📌 Primitive Type: Kiểu dữ liệu cơ bản, lưu giá trị trực tiếp

🔢 8 Kiểu Primitive:

Kiểu      Mô tả           Kích thước   Mặc định   Phạm vi
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
byte      Số nguyên nhỏ   8-bit        0          -128 đến 127
short     Số nguyên       16-bit       0          -32,768 đến 32,767
int       Số nguyên       32-bit       0          -2³¹ đến 2³¹-1
long      Số nguyên lớn   64-bit       0L         -2⁶³ đến 2⁶³-1
float     Số thực đơn     32-bit       0.0f       ~±3.4E+38
double    Số thực kép     64-bit       0.0d       ~±1.7E+308
char      Ký tự Unicode   16-bit       '\\u0000'   0 đến 65,535
boolean   true/false      1-bit        false      true/false"""
    add_content_box(slide8, content, 1.3, 5.5)

    # Slide 9: Reference Types
    slide9 = add_slide_with_title("Slide 7.1 – Reference Type (Kiểu tham chiếu)")
    content = """📦 Reference Type: Lưu trữ địa chỉ bộ nhớ, không phải giá trị

Đặc điểm:
• Lưu địa chỉ của object trong Stack
• Object thực tế nằm trong Heap
• Có method và thuộc tính
• Giá trị mặc định là null

🎯 Các loại Reference Type:
1. Class (String, Scanner, Student...)
2. Array (int[], String[]...)
3. Interface (List, Map...)

💻 Ví dụ:
String s1 = "Hello";
String s2 = s1;  // s2 lưu địa chỉ, không phải "Hello"

StringBuilder sb1 = new StringBuilder("Hi");
StringBuilder sb2 = sb1;
sb2.append(" there");
System.out.println(sb1);  // "Hi there" - sb2 ảnh hưởng sb1"""
    add_content_box(slide9, content, 1.3, 5.5)

    # Slide 10: Primitive vs Reference
    slide10 = add_slide_with_title("Slide 7.2 – So sánh Primitive vs Reference")
    content = """📊 Bảng so sánh:

Đặc điểm       Primitive Type        Reference Type
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Lưu trữ        Giá trị trực tiếp     Địa chỉ bộ nhớ
Vị trí         Stack                 Stack + Heap
Mặc định       0, false              null
Kích thước     Cố định               Không cố định
So sánh ==     So sánh giá trị       So sánh địa chỉ
Có method?     ❌ Không              ✅ Có
Có thể null?   ❌ Không              ✅ Có
Tốc độ         ⚡ Nhanh hơn          🐢 Chậm hơn

🎯 Lý do chia 2 loại:
1. Hiệu suất - Primitive cực nhanh
2. Tiết kiệm bộ nhớ - Primitive chỉ lưu giá trị
3. Đơn giản hóa - Dễ dùng, không lo null"""
    add_content_box(slide10, content, 1.3, 5.5)

    # Slide 11: Control Flow
    slide11 = add_slide_with_title("Slide 8 – Câu lệnh điều kiện")
    content = """🔀 Các loại câu lệnh:
1. if – điều kiện đơn
2. if-else – chọn 1 trong 2 trường hợp
3. else if – chuỗi điều kiện
4. switch – lựa chọn theo giá trị cụ thể

💻 Ví dụ:
if (score >= 90) {
    System.out.println("A");
} else if (score >= 80) {
    System.out.println("B");
} else {
    System.out.println("C");
}

switch (day) {
    case "Monday":
        System.out.println("Start of week");
        break;
    default:
        System.out.println("Mid week");
}"""
    add_content_box(slide11, content, 1.3, 5.5)

    # Slide 12: 1D Array
    slide12 = add_slide_with_title("Slide 9 – Mảng 1 chiều (1D Array)")
    content = """📦 Định nghĩa: Tập hợp phần tử cùng kiểu dữ liệu

💻 Khai báo & Khởi tạo:
int[] arr = new int[5];
int[] nums = {1, 2, 3, 4, 5};

🔧 Thuộc tính: arr.length

🔄 Duyệt mảng:
// For loop
for (int i = 0; i < arr.length; i++) {
    System.out.println(arr[i]);
}

// Foreach
for (int num : nums) {
    System.out.println(num);
}

⚠️ Lỗi thường gặp: ArrayIndexOutOfBoundsException"""
    add_content_box(slide12, content, 1.3, 5.5)

    # Slide 13: 2D Array
    slide13 = add_slide_with_title("Slide 10 – Mảng 2 chiều (2D Array)")
    content = """📊 Định nghĩa: Mảng của mảng

💻 Khai báo & Khởi tạo:
int[][] matrix = new int[3][4];
int[][] m = {{1, 2}, {3, 4}};

🎯 Ứng dụng:
• Bảng dữ liệu
• Ma trận
• Grid layout

🔄 Duyệt mảng 2 chiều:
for (int i = 0; i < matrix.length; i++) {
    for (int j = 0; j < matrix[i].length; j++) {
        System.out.print(matrix[i][j] + " ");
    }
    System.out.println();
}"""
    add_content_box(slide13, content, 1.3, 5.5)

    # Slide 14: Methods
    slide14 = add_slide_with_title("Slide 11 – Hàm (Method) trong Java")
    content = """🎯 Mục đích:
• Tái sử dụng logic
• Tách biệt chức năng
• Giảm code lặp

💻 Cú pháp:
returnType methodName(params) {
    // code
    return value;
}

🔄 Method Overloading: Cùng tên, khác tham số
public int sum(int a, int b) { return a + b; }
public double sum(double a, double b) { return a + b; }
public int sum(int a, int b, int c) { return a + b + c; }

✅ Best Practice:
• Mỗi method chỉ làm 1 nhiệm vụ
• Tên dùng camelCase và rõ nghĩa
• Không quá dài (< 20 dòng)"""
    add_content_box(slide14, content, 1.3, 5.5)

    # Slide 15: Summary
    slide15 = add_slide_with_title("Slide 12 – Tổng kết")
    content = """🎓 Java Basic – Nền tảng gồm:

1. ✅ Biến & Kiểu dữ liệu
   • Local, Instance, Static variables
   • Primitive types: int, double, char, boolean

2. ✅ Điều kiện
   • if-else, else if
   • switch-case

3. ✅ Mảng
   • Mảng 1 chiều (1D Array)
   • Mảng 2 chiều (2D Array)

4. ✅ Hàm (Method)
   • Khai báo & sử dụng
   • Overloading, Best practices

🚀 Đây là nền móng để học:
• Selenium WebDriver
• OOP (Object-Oriented Programming)
• Test Automation Framework"""
    add_content_box(slide15, content, 1.3, 5.5)

    # Slide 16: Exercises
    slide16 = add_slide_with_title("Slide 13 – Bài tập thực hành tổng hợp")
    content = """🎯 Mục tiêu: Ứng dụng tất cả kiến thức đã học

📝 BÀI TẬP 1: Quản lý điểm sinh viên
• Tạo mảng lưu điểm 5 sinh viên
• Tính điểm trung bình, tìm max, min
• Đếm số sinh viên đạt/không đạt

📝 BÀI TẬP 2: Calculator với Method Overloading
• Tạo Calculator với static counter
• Overload calculate() cho 2 số, 3 số, double

📝 BÀI TẬP 3: Phân loại học lực
• Nhập điểm 0-10
• Phân loại: Xuất sắc, Giỏi, Khá, TB, Yếu

📝 BÀI TẬP 4: Ma trận và tính tổng
• Tạo ma trận 3x3, tính tổng hàng/cột
• Tìm phần tử lớn nhất

📝 BÀI TẬP 5: Kiểm tra số nguyên tố
• Viết method isPrime()
• Tìm và đếm số nguyên tố từ 1-N"""
    add_content_box(slide16, content, 1.3, 5.5)

    # Save presentation
    prs.save('D:/project/java/training/JAVA_BASIC_TRAINING.pptx')
    print("✅ PowerPoint file created successfully: JAVA_BASIC_TRAINING.pptx")

if __name__ == "__main__":
    create_dark_presentation()

